import os
from pptx import Presentation
from typing import Dict, Any, List, Optional
from .base_processor import BaseProcessor
import logging

logger = logging.getLogger(__name__)

class PPTProcessor(BaseProcessor):
    """PPT文件处理器"""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
    
    def process(self, file_path: str) -> Dict[str, Any]:
        """处理PPT文件，提取内容和识别人物"""
        try:
            self.logger.info(f"开始处理PPT文件: {file_path}")
            
            # 验证文件
            if not self.validate_file(file_path):
                return {
                    "status": "error",
                    "file_path": file_path,
                    "error": "无效的PPT文件"
                }
            
            # 打开PPT文件
            prs = Presentation(file_path)
            
            # 提取内容
            content = self._extract_content(prs)
            
            # 识别人物
            persons = self._identify_persons(content)
            
            # 提取演示文稿信息
            presentation_info = self.extract_presentation_info(file_path)
            
            # 获取文件信息
            file_info = self.get_file_info(file_path)
            
            self.logger.info(f"PPT文件处理完成，识别到 {len(persons)} 个人物")
            
            return {
                "status": "success",
                "file_path": file_path,
                "file_info": file_info,
                "presentation_info": presentation_info,
                "content": content,
                "persons": persons
            }
        except Exception as e:
            self.logger.error(f"处理PPT文件时出错: {str(e)}")
            return {
                "status": "error",
                "file_path": file_path,
                "error": str(e)
            }
    
    def validate_file(self, file_path: str) -> bool:
        """验证PPT文件是否有效"""
        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                return False
            
            # 检查文件扩展名
            ext = os.path.splitext(file_path)[1].lower()
            if ext not in [".pptx", ".ppt"]:
                return False
            
            # 尝试打开文件
            prs = Presentation(file_path)
            return True
        except:
            return False
    
    def _extract_content(self, prs: Presentation) -> Dict[str, Any]:
        """提取PPT内容"""
        content = {
            "slides": [],
            "text": ""
        }
        
        for i, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": i + 1,
                "title": "",
                "content": ""
            }
            
            # 提取标题
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if hasattr(shape, 'text') and shape.text:
                        # 假设第一个文本框是标题
                        if not slide_content['title']:
                            slide_content['title'] = shape.text
                        else:
                            slide_content['content'] += shape.text + "\n"
            
            content["slides"].append(slide_content)
            content["text"] += slide_content['title'] + "\n" + slide_content['content'] + "\n"
        
        return content
    
    def _identify_persons(self, content: Dict[str, Any]) -> List[str]:
        """从内容中识别人物"""
        persons = []
        
        # 简单的人物识别逻辑，实际应用中可以使用更复杂的NLP方法
        text = content.get("text", "")
        
        # 这里可以添加更复杂的人物识别逻辑
        # 例如使用命名实体识别(NER)模型
        
        # 暂时返回一个空列表，后续会实现更复杂的人物识别
        return persons
    
    def extract_presentation_info(self, file_path: str) -> Dict[str, Any]:
        """提取PPT演示文稿信息"""
        try:
            prs = Presentation(file_path)
            
            info = {
                "slide_count": len(prs.slides),
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "created": str(prs.core_properties.created) if prs.core_properties.created else "",
                "modified": str(prs.core_properties.modified) if prs.core_properties.modified else ""
            }
            
            return info
        except Exception as e:
            self.logger.error(f"提取PPT信息时出错: {str(e)}")
            return {}
